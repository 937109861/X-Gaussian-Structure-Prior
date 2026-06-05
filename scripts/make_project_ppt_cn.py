# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = "XGaussian_Project_Progress_Report_CN.pptx"
FONT = "Microsoft YaHei"

COLORS = {
    "navy": RGBColor(28, 55, 85),
    "blue": RGBColor(48, 101, 160),
    "teal": RGBColor(36, 132, 128),
    "green": RGBColor(78, 143, 92),
    "orange": RGBColor(203, 122, 60),
    "red": RGBColor(178, 75, 75),
    "gray": RGBColor(92, 99, 112),
    "light": RGBColor(245, 247, 250),
    "mid": RGBColor(215, 221, 229),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(34, 38, 46),
}


def set_font(paragraph, size=14, bold=False, color=None, align=None):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    if color is not None:
        paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def text_box(slide, x, y, w, h, text, size=14, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    set_font(p, size=size, bold=bold, color=color or COLORS["black"], align=align)
    return box


def title(slide, text, sub=None):
    text_box(slide, 0.55, 0.25, 12.2, 0.55, text, 24, True, COLORS["navy"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.9), Inches(12.2), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.fill.background()
    if sub:
        text_box(slide, 0.58, 0.98, 12.0, 0.28, sub, 10, False, COLORS["gray"])


def footer(slide):
    text_box(slide, 0.55, 7.12, 12.2, 0.25, "X-Gaussian 项目阶段性汇报", 8, False, COLORS["gray"], PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, heading, bullets, accent="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["light"]
    shape.line.color.rgb = COLORS["mid"]
    text_box(slide, x + 0.18, y + 0.12, w - 0.36, 0.35, heading, 15, True, COLORS[accent])
    body = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.58), Inches(w - 0.44), Inches(h - 0.72))
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        set_font(p, size=12, color=COLORS["black"])
        p.space_after = Pt(5)
    return shape


def process(slide, x, y, w, h, text, color="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[color]
    shape.line.color.rgb = COLORS[color]
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    set_font(p, size=10.5, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    return shape


def arrow(slide, x1, y1, x2, y2):
    con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    con.line.color.rgb = COLORS["gray"]
    con.line.width = Pt(1.1)
    con.line.end_arrowhead = True


def metric(slide, x, y, w, h, label, value, note, color="teal"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 251, 253)
    shape.line.color.rgb = COLORS["mid"]
    text_box(slide, x + 0.12, y + 0.12, w - 0.24, 0.28, label, 10.5, True, COLORS["gray"], PP_ALIGN.CENTER)
    text_box(slide, x + 0.08, y + 0.46, w - 0.16, 0.5, value, 21, True, COLORS[color], PP_ALIGN.CENTER)
    text_box(slide, x + 0.12, y + 1.02, w - 0.24, 0.35, note, 8.5, False, COLORS["gray"], PP_ALIGN.CENTER)


def placeholder(slide, x, y, w, h, heading, hint):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = COLORS["mid"]
    shape.line.width = Pt(1.4)
    text_box(slide, x + 0.15, y + 0.12, w - 0.3, 0.35, heading, 13, True, COLORS["navy"], PP_ALIGN.CENTER)
    text_box(slide, x + 0.25, y + h / 2 - 0.28, w - 0.5, 0.85, hint, 11, False, COLORS["gray"], PP_ALIGN.CENTER)


def table(slide, x, y, w, h, data, font_size=10):
    rows, cols = len(data), len(data[0])
    shp = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tab = shp.table
    for r in range(rows):
        for c in range(cols):
            cell = tab.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["navy"]
            elif r % 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 240, 246)
            for p in cell.text_frame.paragraphs:
                set_font(p, size=font_size if r else font_size + 1, bold=(r == 0), color=COLORS["white"] if r == 0 else COLORS["black"], align=PP_ALIGN.CENTER)
    return shp


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(247, 249, 252)
    return s


# 1
s = blank()
text_box(s, 0.85, 1.25, 11.8, 0.8, "X-Gaussian 稀疏视角 CT 重建项目", 34, True, COLORS["navy"], PP_ALIGN.CENTER)
text_box(s, 1.1, 2.18, 11.2, 0.65, "阶段性进展汇报：从二维新视角渲染到三维结构体数据生成", 19, False, COLORS["gray"], PP_ALIGN.CENTER)
metric(s, 1.7, 4.0, 2.25, 1.25, "二维 PSNR", "45.61 dB", "最终方法 30k", "teal")
metric(s, 4.35, 4.0, 2.25, 1.25, "二维 SSIM", "0.9999", "最终方法 30k", "blue")
metric(s, 7.0, 4.0, 2.25, 1.25, "ROI 区域 PSNR", "13.30 dB", "三维体数据", "green")
metric(s, 9.65, 4.0, 2.25, 1.25, "有效结构 PSNR", "14.28 dB", "occupancy 区域", "orange")
text_box(s, 0.9, 6.55, 11.6, 0.35, "汇报人：__________    数据集：chest_50.pickle    当前阶段：单数据集完整流程验证完成", 12, False, COLORS["gray"], PP_ALIGN.CENTER)

# 2
s = blank(); title(s, "1. 研究问题：为什么要改 X-Gaussian？"); footer(s)
card(s, 0.75, 1.35, 3.7, 4.9, "原方法优势", ["可以用稀疏 X-ray 投影训练 Gaussian 表示", "二维新视角渲染速度快、指标较好", "输出 Gaussian 点云，具有一定可解释性"], "blue")
card(s, 4.8, 1.35, 3.7, 4.9, "存在的问题", ["训练主要受二维投影约束", "三维空间中 Gaussian 分布不一定合理", "直接导出 3D Volume 时结构弱、杂质多", "3D Slicer 中容易出现黑、碎、不清晰"], "red")
card(s, 8.85, 1.35, 3.7, 4.9, "本项目目标", ["保持二维渲染质量", "引入三维结构先验", "让 Gaussian 更集中于有效结构区域", "提升三维体数据的结构表达能力"], "teal")

# 3
s = blank(); title(s, "2. 项目总体思路"); footer(s)
process(s, 0.8, 1.35, 2.55, 0.85, "输入数据\n稀疏 X-ray + CT 体数据", "blue")
process(s, 3.8, 1.35, 2.55, 0.85, "构建先验\nROI / occupancy / prior_volume", "teal")
process(s, 6.8, 1.35, 2.55, 0.85, "训练模型\n先验引导 X-Gaussian", "green")
process(s, 9.8, 1.35, 2.55, 0.85, "输出结果\n二维图像 + 三维体数据", "orange")
arrow(s, 3.35, 1.78, 3.8, 1.78); arrow(s, 6.35, 1.78, 6.8, 1.78); arrow(s, 9.35, 1.78, 9.8, 1.78)
card(s, 0.9, 3.0, 5.55, 2.6, "判断标准", ["二维渲染 PSNR / SSIM 不下降", "Gaussian 点云更靠近有效结构", "三维 ROI 和 occupancy 区域指标提升", "背景无效响应降低"], "navy")
card(s, 6.9, 3.0, 5.55, 2.6, "当前完成情况", ["chest 数据集完整流程已经跑通", "已完成 baseline 30k 与 ours 30k 对比", "二维指标高，三维结构区域提升明显", "已具备向老师汇报和继续推进的基础"], "green")

# 4
s = blank(); title(s, "3. 总体流程架构图"); footer(s)
process(s, 0.55, 1.25, 1.7, 0.75, "输入\nchest_50.pickle", "navy")
process(s, 2.65, 0.95, 2.05, 0.75, "稀疏 X-ray 投影\n训练/测试视角", "blue")
process(s, 2.65, 2.05, 2.05, 0.75, "真实 CT 体数据\n构建先验", "blue")
process(s, 5.15, 2.05, 2.15, 0.75, "三维结构先验\nROI / occupancy / prior_volume", "teal")
process(s, 5.15, 0.95, 2.15, 0.75, "ACUI 初始化\nGaussian 点云", "teal")
process(s, 7.85, 1.5, 2.25, 0.78, "先验引导训练\nX-Gaussian", "green")
process(s, 10.65, 0.75, 2.05, 0.75, "二维新视角图像\nrenders", "orange")
process(s, 10.65, 1.85, 2.05, 0.75, "三维点云\npoint_cloud.ply", "orange")
process(s, 7.85, 3.35, 2.25, 0.75, "结构保持体生成\nGaussian-to-Volume", "green")
process(s, 10.65, 3.15, 2.05, 0.75, "三维体数据\nrecon_volume.npy", "orange")
process(s, 10.65, 4.15, 2.05, 0.75, "Slicer 文件\n.mhd / .raw", "orange")
process(s, 5.15, 4.4, 2.15, 0.75, "指标评估\n2D + 3D", "teal")
arrow(s, 2.25, 1.62, 2.65, 1.32); arrow(s, 2.25, 1.62, 2.65, 2.42)
arrow(s, 4.7, 2.42, 5.15, 2.42); arrow(s, 4.7, 1.32, 5.15, 1.32)
arrow(s, 7.3, 1.32, 7.85, 1.85); arrow(s, 7.3, 2.42, 7.85, 1.9)
arrow(s, 10.1, 1.86, 10.65, 1.12); arrow(s, 10.1, 1.9, 10.65, 2.22)
arrow(s, 11.65, 2.6, 9.0, 3.35); arrow(s, 7.3, 2.42, 8.45, 3.35)
arrow(s, 10.1, 3.72, 10.65, 3.52); arrow(s, 10.1, 3.78, 10.65, 4.52)
arrow(s, 10.65, 3.52, 7.3, 4.78)

# 5
s = blank(); title(s, "4. 主创新点：三维先验引导的 Gaussian 优化"); footer(s)
card(s, 0.7, 1.3, 3.75, 5.0, "构建结构先验", ["从 CT volume 中提取有效结构区域", "得到 occupancy_mask 和 roi_mask", "保存 prior_data.npz 供训练和导出使用"], "blue")
card(s, 4.8, 1.3, 3.75, 5.0, "训练阶段使用先验", ["ROI 区域引导 densification", "空区域 Gaussian pruning", "统计 Gaussian 与 ROI / occupancy 的空间关系"], "teal")
card(s, 8.9, 1.3, 3.75, 5.0, "预期效果", ["减少无效区域 Gaussian", "使点云更贴近真实结构分布", "为后续三维体数据生成提供更好的空间基础"], "green")

# 6
s = blank(); title(s, "5. 副创新点一：Occupancy 一致性约束"); footer(s)
process(s, 1.0, 1.35, 2.55, 0.85, "当前 Gaussian\n位置/尺度/不透明度", "blue")
process(s, 4.25, 1.35, 2.55, 0.85, "低分辨率体素化\nGaussian-to-Grid", "teal")
process(s, 7.5, 1.35, 2.55, 0.85, "先验结构网格\noccupancy prior", "teal")
process(s, 5.9, 3.0, 2.6, 0.85, "L_occ\n结构一致性损失", "green")
arrow(s, 3.55, 1.78, 4.25, 1.78); arrow(s, 6.8, 1.78, 7.5, 1.78)
arrow(s, 5.5, 2.2, 6.5, 3.0); arrow(s, 8.15, 2.2, 7.3, 3.0)
card(s, 1.0, 4.35, 11.0, 1.45, "作用说明", ["该模块主要提供三维结构正则，目标不是单独提高二维 PSNR，而是让 Gaussian 分布更接近最终需要的三维结构。", "当前实验显示：该模块没有破坏原模型训练，可作为结构引导项保留。"], "navy")

# 7
s = blank(); title(s, "6. 副创新点二：结构保持的体数据生成"); footer(s)
card(s, 0.7, 1.25, 3.7, 5.0, "原始导出的问题", ["直接使用 opacity 做 splatting", "opacity 不等价于 CT 密度", "有效结构弱，背景杂质多", "Slicer 中显示效果不稳定"], "red")
card(s, 4.8, 1.25, 3.7, 5.0, "改进后的生成策略", ["结合 opacity、feature 和 Gaussian scale", "使用 ROI / occupancy gate", "融合 prior_volume 连续强度", "增强结构区域，压低背景区域"], "teal")
card(s, 8.9, 1.25, 3.7, 5.0, "最终输出", ["recon_volume.npy", "recon_volume.mhd / raw", "recon_volume_slicer.mhd / raw", "volume_metrics.json", "三切面预览图"], "green")

# 8
s = blank(); title(s, "7. 二维新视角渲染结果"); footer(s)
data = [
    ["迭代轮数", "PSNR", "SSIM", "FPS", "Gaussian 数量"],
    ["1000", "30.02", "0.9975", "421.52", "11984"],
    ["2000", "34.00", "0.9989", "347.39", "50955"],
    ["5000", "38.04", "0.9995", "313.26", "72622"],
    ["10000", "42.92", "0.9998", "303.70", "101487"],
    ["20000", "45.55", "0.9999", "294.30", "130909"],
    ["30000", "45.61", "0.9999", "300.66", "130909"],
]
table(s, 0.75, 1.35, 7.65, 4.55, data, 10)
card(s, 8.75, 1.35, 3.75, 4.55, "结果说明", ["训练过程稳定收敛", "20k 后基本达到饱和", "30k 最终 PSNR 为 45.61 dB", "说明二维新视角合成质量较高", "后续重点比较：在保持二维质量的同时提升三维结构"], "green")

# 9
s = blank(); title(s, "8. 三维体数据指标对比"); footer(s)
data = [
    ["方法", "整体 PSNR", "整体 MAE", "ROI PSNR", "有效结构 PSNR", "ROI 外平均响应"],
    ["Baseline 30k", "10.71", "0.2175", "6.37", "3.95", "0.0169"],
    ["Ours 30k", "11.40", "0.1984", "13.30", "14.28", "0.0040"],
    ["提升", "+0.69", "-0.0191", "+6.93", "+10.33", "降低约 76%"],
]
table(s, 0.55, 1.3, 12.25, 2.15, data, 9)
metric(s, 1.05, 4.1, 2.4, 1.25, "ROI 区域", "+6.93 dB", "结构区域明显提升", "green")
metric(s, 3.95, 4.1, 2.4, 1.25, "有效结构区域", "+10.33 dB", "occupancy 区域提升", "orange")
metric(s, 6.85, 4.1, 2.4, 1.25, "背景响应", "-76%", "杂质明显减少", "teal")
metric(s, 9.75, 4.1, 2.4, 1.25, "整体误差", "0.1984", "低于 baseline", "blue")

# 10
s = blank(); title(s, "9. 二维图像可视化对比（待放图）"); footer(s)
placeholder(s, 0.75, 1.35, 3.85, 4.65, "Baseline 渲染图", "放 baseline_30k 的 test/ours_30000/renders 示例图")
placeholder(s, 4.75, 1.35, 3.85, 4.65, "本文方法渲染图", "放 ours_30k 的 test/ours_30000/renders 同编号图")
placeholder(s, 8.75, 1.35, 3.85, 4.65, "真实图 / 误差图", "放 GT 图，或放 absdiff 误差图")

# 11
s = blank(); title(s, "10. 三维体数据可视化对比（待放图）"); footer(s)
placeholder(s, 0.7, 1.2, 3.9, 2.3, "Baseline 3D Slicer", "放 baseline volume 的 3D Slicer 截图")
placeholder(s, 4.75, 1.2, 3.9, 2.3, "本文方法 3D Slicer", "放 structure_density volume 的 3D Slicer 截图")
placeholder(s, 8.8, 1.2, 3.8, 2.3, "参考图 / GT", "放 GT volume 截图或参考三维结构图")
placeholder(s, 0.7, 4.05, 3.9, 2.0, "Baseline 三切面", "放 axial / coronal / sagittal 切片")
placeholder(s, 4.75, 4.05, 3.9, 2.0, "本文方法三切面", "放 ours 的三切面切片")
placeholder(s, 8.8, 4.05, 3.8, 2.0, "误差图", "放 absdiff 或误差热力图")

# 12
s = blank(); title(s, "11. 当前结论与不足"); footer(s)
card(s, 0.75, 1.25, 5.75, 4.9, "当前结论", ["chest 数据集上完整流程已经跑通", "二维新视角渲染质量较高", "三维 ROI 和 occupancy 区域指标明显优于 baseline", "背景无效响应显著降低", "项目具备继续推进为小论文的基础"], "green")
card(s, 6.85, 1.25, 5.75, 4.9, "需要注意的问题", ["目前主要只验证了一个数据集", "prior 来源需要在论文中解释清楚", "整体 Volume PSNR 仍然不高", "Slicer 可视化仍可能存在少量碎片", "还需要补消融实验和更多可视化证据"], "orange")

# 13
s = blank(); title(s, "12. 下一步计划"); footer(s)
card(s, 0.75, 1.2, 3.75, 5.0, "实验补充", ["完成更多数据集或不同病例", "补充不同稀疏视角设置", "整理 baseline 与本文方法的完整二维/三维表格", "补充消融实验"], "blue")
card(s, 4.8, 1.2, 3.75, 5.0, "方法优化", ["加入小连通域过滤，减少孤立杂质", "尝试更合理的 prior 来源，例如 FDK 或 NAF", "优化 Gaussian-to-Volume 参数", "提升 Slicer 可视化质量"], "teal")
card(s, 8.85, 1.2, 3.75, 5.0, "论文准备", ["明确问题、方法、实验的叙事逻辑", "突出结构区域恢复，而非完整软组织 CT 恢复", "整理可视化证据链", "形成小论文初稿框架"], "green")

prs.save(OUT)
print(OUT)
